"""
Markdown presenter-deck to styled PowerPoint converter.

Usage:
    python generate_pptx.py

Reads all presenter-deck.md files from materials/*/slides/ and produces
professional PPTX files in pptx-output/.

Requires: python-pptx (pip install python-pptx)
"""

import re
import os
import glob
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE


# ---------------------------------------------------------------------------
# Color palette -- professional blue / dark theme
# ---------------------------------------------------------------------------
COLOR_PRIMARY = RGBColor(0x1B, 0x3A, 0x5C)      # dark navy
COLOR_ACCENT = RGBColor(0x00, 0x78, 0xD4)        # Microsoft blue
COLOR_ACCENT2 = RGBColor(0x2E, 0xCC, 0x71)       # green accent
COLOR_ACCENT3 = RGBColor(0xE7, 0x4C, 0x3C)       # red/warm accent
COLOR_LIGHT = RGBColor(0xEC, 0xF0, 0xF1)         # light gray bg
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_DARK_TEXT = RGBColor(0x2C, 0x3E, 0x50)
COLOR_MID_TEXT = RGBColor(0x5D, 0x6D, 0x7E)
COLOR_NOTES_BG = RGBColor(0xFD, 0xF2, 0xE9)      # warm hint for notes shapes

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Icon-like shapes for slide categories
ICON_MAP = {
    "title": MSO_SHAPE.ROUNDED_RECTANGLE,
    "objectives": MSO_SHAPE.PENTAGON,
    "example": MSO_SHAPE.CHEVRON,
    "lab": MSO_SHAPE.FLOWCHART_PROCESS,
    "validation": MSO_SHAPE.FLOWCHART_DECISION,
    "reflection": MSO_SHAPE.CLOUD,
    "references": MSO_SHAPE.FOLDED_CORNER,
    "default": MSO_SHAPE.ROUNDED_RECTANGLE,
}


def categorize_slide(title_lower):
    """Return a category key based on slide title keywords."""
    if any(w in title_lower for w in ["title", "module purpose"]):
        return "title"
    if any(w in title_lower for w in ["objective", "learning"]):
        return "objectives"
    if any(w in title_lower for w in ["example", "bi example", "ds example", "de example", "scenario"]):
        return "example"
    if any(w in title_lower for w in ["lab", "hands-on", "exercise"]):
        return "lab"
    if any(w in title_lower for w in ["validation", "checklist"]):
        return "validation"
    if any(w in title_lower for w in ["reflection", "wrap"]):
        return "reflection"
    if any(w in title_lower for w in ["reference", "resource"]):
        return "references"
    return "default"


def accent_for_category(cat):
    """Return accent color for slide category."""
    mapping = {
        "title": COLOR_ACCENT,
        "objectives": COLOR_ACCENT,
        "example": COLOR_ACCENT2,
        "lab": COLOR_ACCENT2,
        "validation": COLOR_ACCENT3,
        "reflection": COLOR_PRIMARY,
        "references": COLOR_MID_TEXT,
        "default": COLOR_ACCENT,
    }
    return mapping.get(cat, COLOR_ACCENT)


# ---------------------------------------------------------------------------
# Markdown parser
# ---------------------------------------------------------------------------

def parse_presenter_deck(md_path):
    """Parse a presenter-deck.md into a list of slide dicts."""
    with open(md_path, "r", encoding="utf-8") as f:
        text = f.read()

    # Extract deck title from H1
    deck_title = ""
    h1_match = re.match(r"^#\s+(.+)", text, re.MULTILINE)
    if h1_match:
        deck_title = h1_match.group(1).strip()
        # Remove "Presenter Deck: " prefix if present
        deck_title = re.sub(r"^Presenter Deck:\s*", "", deck_title)

    # Split on H2 headings (## Slide N. Title  or  ## Slide N: Title)
    slide_blocks = re.split(r"\n(?=## )", text)
    slides = []

    for block in slide_blocks:
        h2_match = re.match(r"^##\s+Slide\s+\d+[.:]\s*(.+)", block)
        if not h2_match:
            continue
        title = h2_match.group(1).strip()
        body = block[h2_match.end():].strip()

        # Parse sections by sub-headings or label patterns
        on_screen = []
        presenter_notes = []
        key_points = []

        current_section = None
        for line in body.split("\n"):
            stripped = line.strip()
            lower = stripped.lower()

            if lower.startswith("on-screen content:") or lower.startswith("key points:"):
                current_section = "content"
                continue
            elif lower.startswith("presenter note") or lower.startswith("presenter notes:"):
                current_section = "notes"
                continue

            if stripped.startswith("- "):
                item = stripped[2:].strip()
                if current_section == "content":
                    on_screen.append(item)
                elif current_section == "notes":
                    presenter_notes.append(item)
                else:
                    on_screen.append(item)

        slides.append({
            "title": title,
            "content": on_screen,
            "notes": presenter_notes,
        })

    return deck_title, slides


# ---------------------------------------------------------------------------
# Shape helpers
# ---------------------------------------------------------------------------

def add_background(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent_bar(slide, color, left=0, top=0, width=Inches(0.15), height=None):
    """Add a vertical accent bar on the left side of the slide."""
    if height is None:
        height = SLIDE_HEIGHT
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_top_band(slide, color, height=Inches(1.4)):
    """Add a colored band across the top."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_WIDTH, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_bottom_strip(slide, color, text, height=Inches(0.45)):
    """Add a footer strip at the bottom with text."""
    top = SLIDE_HEIGHT - height
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, top, SLIDE_WIDTH, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.size = Pt(10)
    run.font.color.rgb = COLOR_WHITE
    return shape


def add_icon_shape(slide, category, left, top, size=Inches(0.6)):
    """Add a small decorative shape as a visual icon."""
    shape_type = ICON_MAP.get(category, ICON_MAP["default"])
    shape = slide.shapes.add_shape(shape_type, left, top, size, size)
    color = accent_for_category(category)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                font_color=COLOR_DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Segoe UI"):
    """Add a simple text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = font_color
    run.font.bold = bold
    run.font.name = font_name
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    font_color=COLOR_DARK_TEXT, bullet_color=None):
    """Add a text box with bulleted items."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(8)
        p.space_before = Pt(4)
        p.level = 0

        # Bullet character
        bullet_run = p.add_run()
        bullet_run.text = "\u25CF  "  # filled circle
        bullet_run.font.size = Pt(font_size - 2)
        bullet_run.font.color.rgb = bullet_color or accent_for_category("default")
        bullet_run.font.name = "Segoe UI"

        # Item text
        run = p.add_run()
        run.text = item
        run.font.size = Pt(font_size)
        run.font.color.rgb = font_color
        run.font.name = "Segoe UI"

    return txBox


def add_notes_box(slide, notes, left, top, width, height):
    """Add a styled presenter-notes box on the slide (subtle callout)."""
    if not notes:
        return
    # Background shape
    bg_shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    bg_shape.fill.solid()
    bg_shape.fill.fore_color.rgb = COLOR_NOTES_BG
    bg_shape.line.color.rgb = RGBColor(0xE0, 0xC8, 0xA8)
    bg_shape.line.width = Pt(1)

    # Label
    label_box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.1), width - Inches(0.4), Inches(0.3)
    )
    ltf = label_box.text_frame
    lp = ltf.paragraphs[0]
    lr = lp.add_run()
    lr.text = "PRESENTER NOTES"
    lr.font.size = Pt(9)
    lr.font.bold = True
    lr.font.color.rgb = COLOR_MID_TEXT
    lr.font.name = "Segoe UI"

    # Notes content
    notes_box = slide.shapes.add_textbox(
        left + Inches(0.2), top + Inches(0.4), width - Inches(0.4), height - Inches(0.5)
    )
    ntf = notes_box.text_frame
    ntf.word_wrap = True
    for i, note in enumerate(notes):
        if i == 0:
            p = ntf.paragraphs[0]
        else:
            p = ntf.add_paragraph()
        p.space_after = Pt(4)
        run = p.add_run()
        run.text = note
        run.font.size = Pt(11)
        run.font.italic = True
        run.font.color.rgb = COLOR_MID_TEXT
        run.font.name = "Segoe UI"


def add_slide_number(slide, number, total):
    """Add slide number in bottom-right."""
    add_textbox(
        slide,
        SLIDE_WIDTH - Inches(1.2), SLIDE_HEIGHT - Inches(0.45),
        Inches(1.0), Inches(0.35),
        f"{number} / {total}",
        font_size=10,
        font_color=COLOR_WHITE,
        alignment=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def build_title_slide(prs, deck_title, subtitle, total_slides):
    """Create a visually rich title slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_background(slide, COLOR_PRIMARY)

    # Large decorative circle (top-right)
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        SLIDE_WIDTH - Inches(4), Inches(-1.5),
        Inches(6), Inches(6),
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x22, 0x4E, 0x7A)
    circle.line.fill.background()

    # Small accent circle
    c2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(0.8), Inches(5.5),
        Inches(1.5), Inches(1.5),
    )
    c2.fill.solid()
    c2.fill.fore_color.rgb = COLOR_ACCENT
    c2.line.fill.background()

    # Title text
    add_textbox(
        slide, Inches(1.0), Inches(2.0), Inches(10), Inches(1.5),
        deck_title, font_size=36, font_color=COLOR_WHITE, bold=True,
    )

    # Subtitle
    add_textbox(
        slide, Inches(1.0), Inches(3.8), Inches(8), Inches(0.8),
        subtitle, font_size=18, font_color=RGBColor(0xBD, 0xC3, 0xC7),
    )

    # Bottom line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(1.0), Inches(5.0), Inches(3), Inches(0.05),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT
    line.line.fill.background()

    # Branding text
    add_textbox(
        slide, Inches(1.0), Inches(5.2), Inches(6), Inches(0.4),
        "AI Coding Agent Training  |  GitHub Copilot in VS Code",
        font_size=12, font_color=RGBColor(0x95, 0xA5, 0xA6),
    )


def build_content_slide(prs, slide_data, slide_num, total_slides, deck_title):
    """Create a styled content slide with accent bar, icon, bullets, and notes."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    add_background(slide, COLOR_WHITE)

    title = slide_data["title"]
    content = slide_data["content"]
    notes = slide_data["notes"]
    category = categorize_slide(title.lower())
    accent = accent_for_category(category)

    # Top accent band
    add_top_band(slide, accent, height=Inches(1.2))

    # Left accent bar (thin)
    add_accent_bar(slide, accent, left=0, top=0, width=Inches(0.08), height=SLIDE_HEIGHT)

    # Slide title on the band
    add_textbox(
        slide, Inches(0.8), Inches(0.25), Inches(10), Inches(0.8),
        title, font_size=28, font_color=COLOR_WHITE, bold=True,
    )

    # Small icon shape
    add_icon_shape(slide, category, Inches(12.0), Inches(0.3), size=Inches(0.6))

    # Slide number badge
    add_textbox(
        slide, Inches(11.8), Inches(0.25), Inches(1.2), Inches(0.4),
        f"{slide_num}", font_size=12, font_color=COLOR_WHITE,
        alignment=PP_ALIGN.RIGHT,
    )

    # Content area
    content_top = Inches(1.6)
    has_notes = len(notes) > 0

    if has_notes:
        content_width = Inches(7.5)
        notes_left = Inches(8.8)
        notes_width = Inches(4.2)
    else:
        content_width = Inches(11.5)

    if content:
        add_bullet_list(
            slide,
            Inches(0.8), content_top, content_width, Inches(4.5),
            content, font_size=18, font_color=COLOR_DARK_TEXT,
            bullet_color=accent,
        )

    # Notes box (right side)
    if has_notes:
        notes_height = min(Inches(1.0 + len(notes) * 0.55), Inches(4.5))
        add_notes_box(slide, notes, notes_left, content_top, notes_width, notes_height)

    # Footer strip
    add_bottom_strip(slide, COLOR_PRIMARY, deck_title)

    # Also add to actual PowerPoint notes field (for presenter view)
    notes_slide = slide.notes_slide
    notes_tf = notes_slide.notes_text_frame
    notes_tf.text = "\n".join(notes) if notes else ""


def build_section_divider(prs, title, subtitle=""):
    """Create a section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLOR_ACCENT)

    # Decorative shapes
    s1 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-2), Inches(-2), Inches(5), Inches(5)
    )
    s1.fill.solid()
    s1.fill.fore_color.rgb = RGBColor(0x00, 0x6A, 0xBE)
    s1.line.fill.background()

    s2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(10), Inches(4), Inches(4), Inches(4)
    )
    s2.fill.solid()
    s2.fill.fore_color.rgb = RGBColor(0x00, 0x6A, 0xBE)
    s2.line.fill.background()

    add_textbox(
        slide, Inches(2), Inches(2.5), Inches(9), Inches(1.5),
        title, font_size=36, font_color=COLOR_WHITE, bold=True,
        alignment=PP_ALIGN.CENTER,
    )
    if subtitle:
        add_textbox(
            slide, Inches(2), Inches(4.2), Inches(9), Inches(0.8),
            subtitle, font_size=18, font_color=RGBColor(0xD6, 0xEA, 0xF8),
            alignment=PP_ALIGN.CENTER,
        )


def build_end_slide(prs, deck_title):
    """Create an end/thank-you slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, COLOR_PRIMARY)

    # Decorative circle
    c = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(4.5), Inches(1.0), Inches(4.3), Inches(4.3)
    )
    c.fill.solid()
    c.fill.fore_color.rgb = COLOR_ACCENT
    c.line.fill.background()

    add_textbox(
        slide, Inches(2), Inches(2.2), Inches(9.3), Inches(1.2),
        "Thank You", font_size=44, font_color=COLOR_WHITE, bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    add_textbox(
        slide, Inches(2), Inches(3.6), Inches(9.3), Inches(0.8),
        deck_title, font_size=18, font_color=RGBColor(0xBD, 0xC3, 0xC7),
        alignment=PP_ALIGN.CENTER,
    )

    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.0), Inches(4.8), Inches(3.3), Inches(0.05),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_ACCENT2
    line.line.fill.background()

    add_textbox(
        slide, Inches(2), Inches(5.2), Inches(9.3), Inches(0.5),
        "AI Coding Agent Training  |  Questions & Discussion",
        font_size=14, font_color=RGBColor(0x95, 0xA5, 0xA6),
        alignment=PP_ALIGN.CENTER,
    )


# ---------------------------------------------------------------------------
# Main pipeline
# ---------------------------------------------------------------------------

def generate_pptx(md_path, output_path):
    """Convert a single presenter-deck.md to a styled PPTX."""
    deck_title, slides = parse_presenter_deck(md_path)
    if not slides:
        print(f"  SKIP (no slides found): {md_path}")
        return

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    total = len(slides)
    subtitle = f"{total} slides  |  BI, Data Science, Data Engineering"

    # Title slide
    build_title_slide(prs, deck_title, subtitle, total)

    # Content slides
    for i, sd in enumerate(slides, 1):
        cat = categorize_slide(sd["title"].lower())

        # Add section dividers before key transitions
        if cat == "lab":
            build_section_divider(prs, "Hands-On Lab", "Time to practice")
        elif cat == "reflection":
            build_section_divider(prs, "Reflection", "What did you learn?")

        build_content_slide(prs, sd, i, total, deck_title)

    # End slide
    build_end_slide(prs, deck_title)

    prs.save(output_path)
    size_kb = os.path.getsize(output_path) / 1024
    print(f"  OK  {os.path.basename(output_path)} ({size_kb:.1f} KB, {total + 3} slides)")


def main():
    base = Path(__file__).parent
    output_dir = base / "pptx-output"
    output_dir.mkdir(exist_ok=True)

    md_files = sorted(glob.glob(
        str(base / "materials" / "*" / "slides" / "presenter-deck.md")
    ))

    if not md_files:
        print("No presenter-deck.md files found.")
        return

    print(f"Found {len(md_files)} presenter decks. Generating styled PPTX...\n")

    for md_path in md_files:
        module_name = Path(md_path).parts[-3]
        out_path = output_dir / f"{module_name}.pptx"
        print(f"  [{module_name}]")
        generate_pptx(md_path, str(out_path))

    print(f"\nDone. Output in: {output_dir}")


if __name__ == "__main__":
    main()
